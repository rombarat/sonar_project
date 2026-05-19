"""
Poster builder — v5.
- Continuous prose everywhere, no bullet points.
- Beam pattern fills the left column.
- Generous gaps between every element.
- Array Setup section removed.
- Bibliography at bottom of middle column.
- Larger body text (22-24 pt).
"""
from __future__ import annotations
import os, sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image as PILImage

HERE     = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = r"C:/Users/akiva/Downloads/PosterA045R.pptx"
OUT      = os.path.join(HERE, "PosterA045R_filled.pptx")
FIGS     = os.path.join(HERE, "figures")

NAVY   = RGBColor(0x0E, 0x28, 0x41)
TEAL   = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
GREEN  = RGBColor(0x19, 0x6B, 0x24)
GRAY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE9, 0xEE, 0xF4)


# ─────────────────────────────────────────────── helpers

def remove_shape(sh):
    el = sh._element; el.getparent().remove(el)

def _run(r, *, name="Calibri", size, bold=False, color):
    f = r.font; f.name = name; f.size = size; f.bold = bold; f.color.rgb = color

def set_text(tf, text, *, sz=Pt(22), bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _run(r, name=font, size=sz, bold=bold, color=color)

def add_para(tf, text, *, sz=Pt(22), bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, before=Pt(0), font="Calibri"):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = before
    r = p.add_run(); r.text = text
    _run(r, name=font, size=sz, bold=bold, color=color)

def tb(slide, L, T, W, H, *, fill=None, border=None):
    sh = slide.shapes.add_textbox(L, T, W, H)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right  = Inches(0.12)
    tf.margin_top  = Inches(0.08); tf.margin_bottom = Inches(0.08)
    if fill:   sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1.2)
    else:      sh.line.fill.background()
    return sh, tf

def sec(slide, L, T, W, text, color=NAVY):
    """34 pt section header + colour underline bar."""
    _, tf = tb(slide, L, T, W, Inches(0.65))
    set_text(tf, text, sz=Pt(34), bold=True, color=color)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, L, T + Inches(0.60), Inches(2.0), Emu(63500))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def sub(slide, L, T, W, text, color=TEAL):
    """24 pt sub-section label."""
    _, tf = tb(slide, L, T, W, Inches(0.42))
    set_text(tf, text, sz=Pt(24), bold=True, color=color)

def fig(slide, fname, L, T, W):
    """Insert figure at width W; return bottom y in inches."""
    path = os.path.join(FIGS, fname)
    im   = PILImage.open(path); pw, ph = im.size
    H    = int(W * ph / pw)
    slide.shapes.add_picture(path, L, T, width=W)
    return T / 914400 + (H / 914400)   # bottom in inches

def cap(slide, L, T, W, text):
    """18 pt italic-style caption."""
    _, tf = tb(slide, L, T, W, Inches(0.52))
    set_text(tf, text, sz=Pt(18), color=GRAY)

def stat(slide, L, T, W, H, number, label, num_color=TEAL):
    _, tf = tb(slide, L, T, W, H, fill=LIGHT)
    set_text(tf, number, sz=Pt(54), bold=True, color=num_color,
             align=PP_ALIGN.CENTER)
    add_para(tf, label, sz=Pt(18), color=GRAY, align=PP_ALIGN.CENTER,
             before=Pt(2))

# shorthand: convert inches to EMU
def I(x): return Inches(x)


# ─────────────────────────────────────────────── build

def build():
    shutil.copyfile(TEMPLATE, OUT)
    prs = Presentation(OUT)
    slide = prs.slides[0]
    by = {sh.name: sh for sh in slide.shapes}

    # ── title bar ─────────────────────────────────────
    tf = by["TextBox 13"].text_frame
    tf.margin_top = I(0.02); tf.margin_bottom = I(0.02)
    set_text(tf, "Hearing the Quiet Vessel",
             sz=Pt(76), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tf = by["TextBox 16"].text_frame
    tf.margin_top = I(0.02); tf.margin_bottom = I(0.02)
    set_text(tf,
        "A Dual-Branch CNN for Sub-Sidelobe Target Detection in Passive Sonar",
        sz=Pt(30), color=WHITE, align=PP_ALIGN.CENTER)
    add_para(tf, "[Student 1], [Student 2]",
             sz=Pt(25), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, before=Pt(6))
    add_para(tf, "Advisor: [Advisor Name, Title]",
             sz=Pt(21), color=WHITE, align=PP_ALIGN.CENTER, before=Pt(2))

    tf = by["TextBox 42"].text_frame
    set_text(tf,
        "Project Number: [XX-X-X-XXXX]   ·   [Institution / Lab]",
        sz=Pt(22), color=WHITE, align=PP_ALIGN.CENTER)

    for nm in ["TextBox 30", "TextBox 19", "TextBox 41", "TextBox 22",
               "Group 1", "תמונה 38", "תמונה 36", "Object 39"]:
        if nm in by: remove_shape(by[nm])

    # ═══════════════════════════════════════════════════════
    # LEFT COLUMN    L=1.27   W=11.49    body 6.19–25.74
    # ═══════════════════════════════════════════════════════
    Lx, Lw = I(1.27), I(11.49)

    # ── Introduction ──────────────────────────────────────
    sec(slide, Lx, I(6.28), Lw, "Introduction")
    y = 7.08

    _, tf = tb(slide, Lx, I(y), Lw, I(3.90))
    set_text(tf,
        "Passive sonar arrays localise underwater acoustic sources through "
        "Delay-and-Sum (DAS) beamforming, a technique that steers a virtual "
        "listening beam by introducing time delays across array elements "
        "proportional to their positions relative to the incoming wavefront.",
        sz=Pt(22))
    add_para(tf,
        "A fundamental limitation arises when multiple sources with disparate "
        "intensities are present simultaneously. The beam pattern of a "
        "100-element uniform linear array (ULA) exhibits a first sidelobe only "
        "13 dB below the main lobe. When a loud vessel dominates the scene, its "
        "sidelobes radiate energy into surrounding bearings and can exceed the "
        "direct contribution of a quieter vessel by 10 dB or more, causing "
        "a conventional energy-threshold detector to miss the quieter target entirely.",
        sz=Pt(22), before=Pt(12))
    add_para(tf,
        "This project proposes a dual-branch convolutional neural network (CNN) "
        "that addresses this sub-sidelobe detection problem by jointly exploiting "
        "the focused-beam power spectrum and the full angular frequency heatmap, "
        "learning to identify the faint target ridge against the dominant "
        "sidelobe background.",
        sz=Pt(22), before=Pt(12))
    y += 3.90 + 0.60

    # beam pattern figure — visual anchor for the problem statement
    bot = fig(slide, "beam_pattern.png", Lx + I(1.25), I(y), I(9.0))
    y = bot + 0.35
    cap(slide, Lx, I(y), Lw,
        "Fig. 1.  Beam pattern of a 100-element ULA steered to 0°. "
        "The quiet vessel (orange marker) is hidden under the −13 dB first "
        "sidelobe of the loud vessel, making it invisible to a simple "
        "energy detector.")
    y += 0.52 + 0.65

    # ── Objectives ────────────────────────────────────────
    sec(slide, Lx, I(y), Lw, "Objectives")
    y += 0.65 + 0.25

    _, tf = tb(slide, Lx, I(y), Lw, I(6.20))
    set_text(tf,
        "The system is evaluated against three quantitative performance targets. "
        "The primary objective (T1) is to detect a quiet vessel when the "
        "interferer's sidelobe at the target direction is at least 10 dB stronger "
        "than the target's direct contribution. This corresponds physically to a "
        "Signal-to-Interference Ratio of at least 23 dB, a demanding regime in "
        "which the sidelobe alone can completely mask the quiet vessel. The v3 "
        "dataset was designed so that 58 % of two-target samples fall within "
        "this T1 regime, ensuring the evaluation reflects realistic sub-sidelobe "
        "detection difficulty.",
        sz=Pt(22))
    add_para(tf,
        "The second objective (T2) requires the network to reliably distinguish "
        "between one and two vessels when their intensities differ by only 3–5 dB. "
        "In this regime, the targets are nearly equal in perceived loudness and "
        "the beamformed output alone provides little discriminating power, "
        "making spectral and spatial cue fusion essential.",
        sz=Pt(22), before=Pt(12))
    add_para(tf,
        "The third objective (T3) extends the problem to source separation, "
        "requiring that distinct audio channels be produced for two overlapping "
        "targets whose intensities differ by no more than 3 dB. "
        "Together, T1–T3 span the continuum from gross detection to fine-grained "
        "source separation in acoustically challenging multi-vessel scenarios.",
        sz=Pt(22), before=Pt(12))

    # ═══════════════════════════════════════════════════════
    # MIDDLE COLUMN    L=13.76   W=11.49    body 6.22–25.74
    # Figures (at W=10.5"):
    #   pipeline_flow  3.837 ratio → h = 2.74"
    #   dualnet_arch   2.030 ratio → h = 5.17"
    # ═══════════════════════════════════════════════════════
    Mx, Mw = I(13.76), I(11.49)
    FW = I(10.50)
    FX = Mx + I(0.50)

    sec(slide, Mx, I(6.30), Mw, "Methods")
    y = 7.10

    sub(slide, Mx, I(y), Mw, "Dataset Generator  (sonar_simulator.py)")
    y += 0.42 + 0.18

    _, tf = tb(slide, Mx, I(y), Mw, I(2.80))
    set_text(tf,
        "A total of 1 000 synthetic samples were generated per dataset version "
        "using sonar_simulator.py. Each sample represents a 3-second acoustic "
        "scene containing one or two vessels. In two-target scenes the quiet "
        "vessel is placed within the main lobe or first sidelobe of the loud "
        "vessel with 80 % probability, constituting the hard detection regime, "
        "and at a well-separated bearing with 20 % probability as the easy "
        "baseline. Source intensities are drawn from the Signal-to-Interference "
        "Ratio range specified for each dataset version.",
        sz=Pt(21))
    add_para(tf,
        "Two complementary representations are computed per sample via a "
        "vectorised Welch sweep: psd_1d, a 513-bin focused-beam power spectral "
        "density at the estimated target bearing, and psd_2d, a full 181 × 513 "
        "angular heatmap. In v3 all samples share a global fixed-range "
        "dB → [0, 1] normalisation, preserving true inter-sample amplitude "
        "differences and preventing the model from exploiting within-sample "
        "normalisation artefacts.",
        sz=Pt(21), before=Pt(10))
    y += 2.80 + 0.60

    bot = fig(slide, "pipeline_flow.png", FX, I(y), FW)
    y = bot + 0.35
    cap(slide, Mx, I(y), Mw,
        "Fig. 2.  Per-sample simulation pipeline. "
        "Hard cases (80 %) position the quiet vessel inside the main lobe "
        "or first sidelobe of the louder interferer.")
    y += 0.52 + 0.65

    sub(slide, Mx, I(y), Mw, "DualBranchNet Architecture", color=ORANGE)
    y += 0.42 + 0.18

    _, tf = tb(slide, Mx, I(y), Mw, I(2.50))
    set_text(tf,
        "The dual-branch architecture processes psd_1d and psd_2d in parallel "
        "and fuses the resulting feature vectors through a shared classification "
        "head, producing a single detection logit that indicates whether the "
        "scene contains one or two vessels.",
        sz=Pt(21))
    add_para(tf,
        "Branch A (Local Focus) applies three Conv1d blocks with progressively "
        "increasing channel depth (1 → 32 → 64 → 128, kernels k = 7 / 5 / 3), "
        "each followed by batch normalisation, ReLU activation, and max-pooling. "
        "Adaptive average pooling yields a 1 024-d feature vector tuned to "
        "narrow tonal signatures in the focused beam.",
        sz=Pt(21), before=Pt(10))
    add_para(tf,
        "Branch B (Global Context) applies four VGG-style double-convolution "
        "blocks (1 → 32 → 64 → 128 → 256 channels) with 2×2 max-pooling to "
        "psd_2d, followed by global average pooling to a 256-d vector encoding "
        "the full sidelobe angular pattern. The 1 280-d concatenation is "
        "passed through three linear layers (1 280 → 256 → 64 → 1) with "
        "ReLU and Dropout (p = 0.3). Total: 1 553 761 parameters.",
        sz=Pt(21), before=Pt(10))
    y += 2.50 + 0.60

    bot = fig(slide, "dualnet_arch.png", FX, I(y), FW)
    y = bot + 0.45

    # bibliography at the bottom of the middle column (compact single block)
    _, tf = tb(slide, Mx, I(y), Mw, I(0.70))
    set_text(tf, "References", sz=Pt(22), bold=True, color=NAVY)
    add_para(tf,
        '[1] Van Trees, Optimum Array Processing, 2002.  '
        '[2] Welch, IEEE Trans. Audio Electroacoust., 1967.  '
        '[3] Simonyan & Zisserman, ICLR, 2015.',
        sz=Pt(18), color=GRAY, before=Pt(5))

    # ═══════════════════════════════════════════════════════
    # RIGHT COLUMN    L=26.61   W=11.49    body 6.19–25.74
    # Figures (widths):
    #   training_curves  2.408 ratio → at 10.8"  h = 4.49"
    #   confusion_matrix 1.129 ratio → at  5.2"  h = 4.61"
    # ═══════════════════════════════════════════════════════
    Rx, Rw = I(26.61), I(11.49)

    sec(slide, Rx, I(6.25), Rw, "Results")
    y = 7.10

    # stat boxes
    sw = I(2.68); sh_ = I(1.80); sgap = (Rw - 4*sw) / 3
    for k, (num, lbl, col) in enumerate([
        ("95.0 %", "Accuracy",  TEAL),
        ("96.8 %", "Precision", TEAL),
        ("92.9 %", "Recall",    ORANGE),
        ("94.8 %", "F1 Score",  TEAL),
    ]):
        stat(slide, Rx + k*(sw+sgap), I(y), sw, sh_, num, lbl, col)
    y += 1.80 + 0.45

    _, tf = tb(slide, Rx, I(y), Rw, I(1.20))
    set_text(tf,
        "The model was trained with AdamW (lr = 10⁻³, weight decay = 10⁻⁴) "
        "and BCEWithLogitsLoss, with the learning rate halved whenever "
        "validation loss plateaued. The best checkpoint was saved at epoch 5 "
        "(val loss = 0.138). Training and validation dynamics are shown below.",
        sz=Pt(21))
    y += 1.20 + 0.50

    bot = fig(slide, "training_curves.png", Rx + I(0.35), I(y), I(10.8))
    y = bot + 0.35
    cap(slide, Rx, I(y), Rw,
        "Fig. 3.  Training dynamics on the v3 dataset. "
        "Validation accuracy peaked at 95.0 % at epoch 5; subsequent "
        "divergence (train → 99.6 %, val ≈ 95 %) indicates overfitting "
        "on the 1 000-sample training set.")
    y += 0.52 + 0.50

    # confusion matrix + annotation side by side
    cm_w = I(5.20)
    bot_cm = fig(slide, "confusion_matrix.png", Rx + I(0.1), I(y), cm_w)

    _, tf = tb(slide, Rx + cm_w + I(0.35), I(y),
               Rw - cm_w - I(0.35), I(4.60))
    set_text(tf,
        "The confusion matrix records 3 false alarms and 7 missed detections "
        "on the 200-sample balanced validation set.",
        sz=Pt(21), bold=False, color=NAVY)
    add_para(tf,
        "The 7 missed detections correspond to the hardest physical regime: "
        "SIR ≥ 30 dB combined with element SNR = −5 dB, placing the quiet "
        "vessel at an output SNR of −20 dB — below the thermal noise floor. "
        "These cases are physically marginal and are not expected to be "
        "resolved without additional data or a larger array.",
        sz=Pt(20), color=GRAY, before=Pt(10))
    add_para(tf,
        "The T1 objective is satisfied: 92.9 % recall was achieved on a "
        "dataset in which 58 % of two-target scenes exercise the SIR ≥ 23 dB "
        "sub-sidelobe regime.",
        sz=Pt(20), color=GREEN, before=Pt(10))

    y = max(bot_cm, (I(y) / 914400) + 4.60) + 0.40

    # ── Conclusions ──────────────────────────────────────
    sec(slide, Rx, I(y), Rw, "Conclusions")
    y += 0.65 + 0.22

    _, tf = tb(slide, Rx, I(y), Rw, I(2.80))
    set_text(tf,
        "DualBranchNet achieves 92.9 % recall on the physics-honest v3 dataset "
        "where 58 % of two-target scenes require sub-sidelobe detection at "
        "SIR ≥ 23 dB, demonstrating that joint spectral and angular fusion "
        "succeeds where energy thresholding fails.",
        sz=Pt(21))
    add_para(tf,
        "Global fixed-range dB normalisation is essential: per-sample "
        "Z-score standardisation destroys inter-sample amplitude contrast, "
        "inflating v1/v2 validation accuracy to a trivial 100 % with no "
        "genuine generalisation.",
        sz=Pt(21), before=Pt(11))
    add_para(tf,
        "Larger training datasets and augmentation are the key next step "
        "to close the train–validation gap (99.6 % vs. 95.0 %) and push "
        "recall toward 100 %.",
        sz=Pt(21), before=Pt(11))

    prs.save(OUT)
    print(f"saved  {OUT}")


if __name__ == "__main__":
    build()
